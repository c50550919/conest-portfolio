#!/usr/bin/env python3
"""
CoNest Pitch Deck Generator
"It Takes a Village" - Safe Housing for Single Parents
$20K Competition Pitch Deck
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Brand colors - From CoNest Investor Deck 2
PRIMARY_GREEN = RGBColor(0x2E, 0x7D, 0x32)  # #2E7D32 - Primary brand color
DARK_GREEN = RGBColor(0x45, 0x7B, 0x3B)  # #457B3B - Headers/accent
ACCENT_BLUE = RGBColor(0x19, 0x76, 0xD2)  # #1976D2 - Blue accent
HIGHLIGHT_ORANGE = RGBColor(0xF5, 0x7C, 0x00)  # #F57C00 - Orange highlight
WARM_CREAM = RGBColor(0xFF, 0xFC, 0xE7)  # #FFFCE7 - Background
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)  # #E8F5E9 - Light green tint
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)  # #E3F2FD - Light blue tint
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)  # #1A1A1A - Text
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)  # #F5F5F5 - Cards
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Logo path
LOGO_PATH = "/Users/ghostmac/Development/conest/workspace/pitch-deck/conest-logo.png"

def set_slide_background(slide, color):
    """Set solid background color for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_shape(slide, text, top=0.5, font_size=44, color=PRIMARY_GREEN):
    """Add a title text box"""
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(1)

    shape = slide.shapes.add_textbox(left, Inches(top), width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    return shape

def add_subtitle(slide, text, top=1.3, font_size=24, color=NEAR_BLACK):
    """Add subtitle text"""
    shape = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.6))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.italic = True
    return shape

def add_body_text(slide, text, top=2.0, left=0.5, width=9, font_size=18):
    """Add body text"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = shape.text_frame
    tf.word_wrap = True

    lines = text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = NEAR_BLACK
        p.space_after = Pt(12)
    return shape

def add_bullet_points(slide, points, top=2.0, left=0.5, font_size=20):
    """Add bullet point list"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(8.5), Inches(5))
    tf = shape.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = NEAR_BLACK
        p.space_after = Pt(16)
        p.level = 0
    return shape

def add_stat_box(slide, number, label, left, top, color=PRIMARY_GREEN):
    """Add a statistics highlight box"""
    # Number
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(2.2), Inches(0.8))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    # Label
    label_shape = slide.shapes.add_textbox(Inches(left), Inches(top + 0.7), Inches(2.2), Inches(0.5))
    tf2 = label_shape.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = label
    p2.font.size = Pt(14)
    p2.font.color.rgb = NEAR_BLACK
    p2.alignment = PP_ALIGN.CENTER

def add_table(slide, headers, rows, top=2.5, left=0.5):
    """Add a table to slide"""
    num_cols = len(headers)
    num_rows = len(rows) + 1

    table = slide.shapes.add_table(num_rows, num_cols, Inches(left), Inches(top),
                                    Inches(9), Inches(num_rows * 0.5)).table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_GREEN
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = NEAR_BLACK

    return table

def create_pitch_deck():
    """Create the full 12-slide pitch deck"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Use blank layout for all slides
    blank_layout = prs.slide_layouts[6]

    # ========== SLIDE 1: Title / Hook ==========
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, WARM_CREAM)

    # Add logo at the top
    import os
    if os.path.exists(LOGO_PATH):
        slide1.shapes.add_picture(LOGO_PATH, Inches(4), Inches(0.3), height=Inches(1.5))

    # Hook line
    hook = slide1.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.5))
    tf = hook.text_frame
    p = tf.paragraphs[0]
    p.text = "Single parents pay 40-60% of their income on rent."
    p.font.size = Pt(28)
    p.font.color.rgb = NEAR_BLACK
    p.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline = slide1.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(1))
    tf2 = tagline.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "We're bringing back the village."
    p2.font.size = Pt(44)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY_GREEN
    p2.alignment = PP_ALIGN.CENTER

    # Company name
    name = slide1.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
    tf3 = name.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "CoNest"
    p3.font.size = Pt(56)
    p3.font.bold = True
    p3.font.color.rgb = DARK_GREEN
    p3.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub = slide1.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(0.5))
    tf4 = sub.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Safe Housing Connections for Single Parents"
    p4.font.size = Pt(20)
    p4.font.color.rgb = NEAR_BLACK
    p4.alignment = PP_ALIGN.CENTER

    # Badge
    badge = slide1.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.4))
    tf5 = badge.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Service-Disabled Veteran-Owned Small Business"
    p5.font.size = Pt(14)
    p5.font.bold = True
    p5.font.color.rgb = HIGHLIGHT_ORANGE
    p5.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 2: The Problem ==========
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, WARM_CREAM)

    add_title_shape(slide2, "The Problem", color=HIGHLIGHT_ORANGE)

    # Story
    story = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.2))
    tf = story.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '"Maria is a nurse making $52,000/year in Charlotte. Her rent is $1,800/month — that\'s 42% of her income. She can\'t save, she can\'t get ahead, and her daughter sees her stress every day."'
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = NEAR_BLACK

    # Stats in boxes
    add_stat_box(slide2, "10.9M", "single-parent households", 0.5, 3.0, PRIMARY_GREEN)
    add_stat_box(slide2, "40-60%", "of income on rent", 3.0, 3.0, HIGHLIGHT_ORANGE)
    add_stat_box(slide2, "35%", "housing cost-burdened", 5.5, 3.0, PRIMARY_GREEN)
    add_stat_box(slide2, "2+ years", "affordable housing waitlist", 8.0, 3.0, HIGHLIGHT_ORANGE)

    # Key message
    key = slide2.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf2 = key.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Housing is unaffordable for single parents who have no safe way to share costs."
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY_GREEN
    p2.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 3: The Solution ==========
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, WARM_CREAM)

    add_title_shape(slide3, "The Solution")
    add_subtitle(slide3, "CoNest makes co-living safe, verified, and matched for compatibility.")

    points = [
        "Mandatory ID verification + background checks (not optional)",
        "Parent-only platform (zero child data stored)",
        "Compatibility matching (schedules, parenting styles, house rules)",
        "End-to-end encrypted messaging"
    ]
    add_bullet_points(slide3, points, top=2.2)

    # Simple math box - Using conservative $600/month (accounts for utilities, shared costs)
    math = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.2), Inches(8), Inches(1))
    math.fill.solid()
    math.fill.fore_color.rgb = PRIMARY_GREEN
    math.line.fill.background()

    math_text = slide3.shapes.add_textbox(Inches(1), Inches(5.4), Inches(8), Inches(0.6))
    tf = math_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Average savings: $600/month per family = $7,200/year back in your pocket"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 4: Value Proposition ==========
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, WARM_CREAM)

    add_title_shape(slide4, "Why CoNest?")

    # Unlike statement - Option B: Emphasize mandatory safety
    unlike = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = unlike.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '"On other platforms, background checks are optional extras most users skip. On CoNest, every parent you meet has been verified."'
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = NEAR_BLACK

    # Comparison table - Fixed semantic accuracy
    headers = ["CoNest", "Traditional Platforms"]
    rows = [
        ["✓ Verification mandatory", "✗ Optional or unavailable"],
        ["✓ ID + background in one $39 fee", "✗ Separate charges (if offered)"],
        ["✓ Parent-focused matching", "✗ Generic filters"],
        ["✓ Zero child data stored", "✗ Mixed demographics"]
    ]
    add_table(slide4, headers, rows, top=2.4)

    # Option C: Cost comparison callout box
    cost_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.0), Inches(9), Inches(1.2))
    cost_box.fill.solid()
    cost_box.fill.fore_color.rgb = LIGHT_GREEN
    cost_box.line.color.rgb = PRIMARY_GREEN

    cost_text = slide4.shapes.add_textbox(Inches(0.6), Inches(5.15), Inches(8.8), Inches(1))
    tf2 = cost_text.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Cost to get verified and message safely:"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = NEAR_BLACK
    p2b = tf2.add_paragraph()
    p2b.text = "SpareRoom: $180+/yr (subscription) + $40 (background) = $220+  •  Roomster: $360+/yr + $30 = $390+  •  CoNest: $39 one-time"
    p2b.font.size = Pt(13)
    p2b.font.color.rgb = NEAR_BLACK

    # ========== SLIDE 5: How It Works ==========
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, WARM_CREAM)

    add_title_shape(slide5, "How It Works")
    add_subtitle(slide5, "Three simple steps to safe, affordable housing")

    # Step 1
    s1 = slide5.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(3), Inches(2))
    tf1 = s1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "1. Sign Up & Get Verified"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY_GREEN
    p1b = tf1.add_paragraph()
    p1b.text = "$39 one-time\nID verification + background check\nTakes <24 hours"
    p1b.font.size = Pt(16)
    p1b.font.color.rgb = NEAR_BLACK

    # Step 2
    s2 = slide5.shapes.add_textbox(Inches(3.5), Inches(2.2), Inches(3), Inches(2))
    tf2 = s2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "2. Match with Parents"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY_GREEN
    p2b = tf2.add_paragraph()
    p2b.text = "Algorithm matches:\n• Schedules (25%)\n• Parenting style (20%)\n• House rules (20%)"
    p2b.font.size = Pt(16)
    p2b.font.color.rgb = NEAR_BLACK

    # Step 3
    s3 = slide5.shapes.add_textbox(Inches(6.5), Inches(2.2), Inches(3), Inches(2))
    tf3 = s3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "3. Move In & Save"
    p3.font.size = Pt(22)
    p3.font.bold = True
    p3.font.color.rgb = PRIMARY_GREEN
    p3b = tf3.add_paragraph()
    p3b.text = "Save $600+/month\nSecure messaging\nShared calendar\nExpense splitting"
    p3b.font.size = Pt(16)
    p3b.font.color.rgb = NEAR_BLACK

    # ========== SLIDE 6: Market Analysis ==========
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, WARM_CREAM)

    add_title_shape(slide6, "Market Opportunity")

    # TAM/SAM/SOM
    headers = ["Segment", "Size", "Calculation"]
    rows = [
        ["TAM", "10.9M households", "All single-parent households in US"],
        ["SAM", "3.8M households", "Housing cost-burdened (35% of TAM)"],
        ["SOM (Year 1)", "1,000 users", "Charlotte metro (~50K single parents)"],
        ["SOM (Year 3)", "20,000 users", "Charlotte + Atlanta + Raleigh-Durham"]
    ]
    add_table(slide6, headers, rows, top=1.8)

    # Launch strategy
    strategy = slide6.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
    tf = strategy.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Launch Strategy:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    for step in ["1. Charlotte first — Prove model with 500+ users",
                 "2. Atlanta second — Expand once Charlotte hits density threshold",
                 "3. Regional expansion — Raleigh-Durham, Nashville (Year 2)"]:
        ps = tf.add_paragraph()
        ps.text = step
        ps.font.size = Pt(16)
        ps.font.color.rgb = NEAR_BLACK

    # ========== SLIDE 7: Financial Model ==========
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, WARM_CREAM)

    add_title_shape(slide7, "Financial Model")

    # Pricing table
    headers = ["Tier", "Price", "What's Included"]
    rows = [
        ["Free", "$0", "Browse profiles"],
        ["Verified", "$39 (one-time)", "ID + background check + messaging"],
        ["Premium", "$14.99/month", "Unlimited swipes, advanced filters"],
        ["Bundle", "$99", "Verification + 6 months premium"]
    ]
    add_table(slide7, headers, rows, top=1.5, left=0.3)

    # Unit economics highlight
    stats = slide7.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2.5))
    tf = stats.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Unit Economics (per 1,000 users)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    metrics = [
        "• Revenue: $33,594  |  Costs: $15,638  |  Profit: $17,956",
        "• 53% overall margin, break-even Month 1",
        "• Cost per verification: $29.89 (Veriff $1.39 + Certn $28.50)",
        "• $39 verification = $7.68 profit after Stripe fees"
    ]
    for m in metrics:
        pm = tf.add_paragraph()
        pm.text = m
        pm.font.size = Pt(16)
        pm.font.color.rgb = NEAR_BLACK

    # ========== SLIDE 8: Traction ==========
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, WARM_CREAM)

    add_title_shape(slide8, "Traction & Progress")

    # Then vs Now table
    headers = ["Then (6 months ago)", "Now"]
    rows = [
        ["Idea on paper", "Functional demo complete"],
        ["No tech stack", "Full mobile apps (iOS + Android)"],
        ["No backend", "Complete API + database"],
        ["No verification plan", "Veriff + Certn contracts signed"],
        ["No payment plan", "Stripe integration built"],
        ["0 waitlist signups", "~500 pre-launch signups"],
        ["No partner discussions", "Housing nonprofits in discussions"]
    ]
    add_table(slide8, headers, rows, top=1.5)

    # What's next
    next_text = slide8.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    tf = next_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Next: Production integrations (4 weeks) → Beta launch → Charlotte public launch"
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY_GREEN
    p.font.bold = True

    # ========== SLIDE 9: Use of $20K ==========
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, WARM_CREAM)

    add_title_shape(slide9, "Use of $20K")
    add_subtitle(slide9, "Current bottleneck: Marketing spend limits user acquisition")

    # Investment table
    headers = ["Investment", "Amount", "Result"]
    rows = [
        ["Social media ads (6 months)", "$9,600", "400+ new users at $24 CAC"],
        ["Partnership outreach", "$4,000", "5 nonprofit partnerships"],
        ["Content marketing", "$3,600", "SEO foundation for organic growth"],
        ["Community events", "$2,800", "10 local meetups"]
    ]
    add_table(slide9, headers, rows, top=2.2)

    # ROI box
    roi = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5), Inches(8), Inches(1.5))
    roi.fill.solid()
    roi.fill.fore_color.rgb = PRIMARY_GREEN
    roi.line.fill.background()

    roi_text = slide9.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1.2))
    tf = roi_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Projected Return: $34,675 revenue from $20K investment"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "73% ROI + 400 families housed"
    p2.font.size = Pt(20)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 10: Team ==========
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, WARM_CREAM)

    add_title_shape(slide10, "Team")

    # Founder
    founder = slide10.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.5))
    tf = founder.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Founder"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN

    for item in ["• Service-Disabled Veteran",
                 "• Full-stack developer",
                 "• Built entire platform solo",
                 "• React Native, Node.js, TypeScript"]:
        pi = tf.add_paragraph()
        pi.text = item
        pi.font.size = Pt(18)
        pi.font.color.rgb = NEAR_BLACK

    # SDVOSB Advantage
    sdvosb = slide10.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(2.5))
    tf2 = sdvosb.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "SDVOSB Advantage"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY_GREEN

    for item in ["• Access to $billions in federal set-asides",
                 "• Sole-source authority up to $5M",
                 "• HUD, VA, DoD contracting opportunities"]:
        pi = tf2.add_paragraph()
        pi.text = item
        pi.font.size = Pt(18)
        pi.font.color.rgb = NEAR_BLACK

    # Advisors
    advisors = slide10.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
    tf3 = advisors.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Advisors Sought: Housing policy expert • Child welfare professional • Government contracting specialist"
    p3.font.size = Pt(16)
    p3.font.color.rgb = NEAR_BLACK

    # ========== SLIDE 11: Call to Action ==========
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, DARK_GREEN)

    # The Ask
    ask = slide11.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1.5))
    tf = ask.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The Ask"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    quote = slide11.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf2 = quote.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = '"Join us in rebuilding the village. With your support, we\'ll house 1,000 families in Year 1 and save each one $7,200 per year."'
    p2.font.size = Pt(24)
    p2.font.italic = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    # What $20K unlocks
    unlocks = slide11.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2.5))
    tf3 = unlocks.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "What $20K Unlocks:"
    p3.font.size = Pt(24)
    p3.font.bold = True
    p3.font.color.rgb = HIGHLIGHT_ORANGE
    p3.alignment = PP_ALIGN.CENTER

    for item in ["• 400+ verified families in Charlotte",
                 "• Proof of product-market fit",
                 "• Expansion trigger for Atlanta",
                 "• Path to seed round"]:
        pi = tf3.add_paragraph()
        pi.text = item
        pi.font.size = Pt(20)
        pi.font.color.rgb = WHITE
        pi.alignment = PP_ALIGN.CENTER

    # Impact math
    impact = slide11.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.6))
    tf4 = impact.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "400 families × $600/month savings = $2.88 million returned to single-parent families annually"
    p4.font.size = Pt(16)
    p4.font.color.rgb = WHITE
    p4.alignment = PP_ALIGN.CENTER

    # ========== SLIDE 12: Closing ==========
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12, PRIMARY_GREEN)

    # Tagline
    tagline = slide12.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.text = "It Takes a Village"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = HIGHLIGHT_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Company name
    name = slide12.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
    tf2 = name.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "CoNest"
    p2.font.size = Pt(64)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub = slide12.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.6))
    tf3 = sub.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Safe Housing Connections for Single Parents"
    p3.font.size = Pt(24)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER

    # Contact placeholder
    contact = slide12.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf4 = contact.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "[Your Name] | [Email] | [Phone]"
    p4.font.size = Pt(18)
    p4.font.color.rgb = WHITE
    p4.alignment = PP_ALIGN.CENTER

    # Website
    web = slide12.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
    tf5 = web.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "conest.com"
    p5.font.size = Pt(20)
    p5.font.bold = True
    p5.font.color.rgb = HIGHLIGHT_ORANGE
    p5.alignment = PP_ALIGN.CENTER

    return prs

if __name__ == "__main__":
    import os

    print("Creating CoNest Pitch Deck...")
    prs = create_pitch_deck()

    # Save to Desktop
    output_path = os.path.expanduser("~/Desktop/CoNest_Pitch_Deck_20K_Competition.pptx")
    prs.save(output_path)
    print(f"✅ Pitch deck saved to: {output_path}")
    print(f"   12 slides created with 'It Takes a Village' narrative")
    print(f"   Ready for $20K competition!")
